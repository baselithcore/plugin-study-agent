"""FastAPI router for the Study Assistant Plugin."""

import io
import json
from typing import Any

from fastapi import APIRouter, File, Form, HTTPException, UploadFile

from core.observability.logging import get_logger
from core.services.llm import get_llm_service
from core.services.voice.models import AudioFormat, VoiceProvider
from core.services.voice.service import VoiceService

from .agent import StudyAgent
from .mcts_engine import OralExamPlanner
from .models import (
    BulkDeleteRequest,
    BulkMoveRequest,
    FlashcardGenerateRequest,
    FlashcardReview,
    OralSessionAnswer,
    OralSessionStart,
    PodcastGenerateRequest,
    SubjectCreate,
)
from .persistence import StudyDAO

logger = get_logger(__name__)


_paddle_ocr_instance = None


def get_paddle_ocr():
    global _paddle_ocr_instance
    if _paddle_ocr_instance is None:
        import sys
        import types

        from langchain_core.documents import Document
        from langchain_text_splitters import RecursiveCharacterTextSplitter

        # Patch langchain imports that paddleocr relies on internally
        if "langchain.docstore.document" not in sys.modules:
            doc_mod = types.ModuleType("langchain.docstore.document")
            doc_mod.Document = Document
            sys.modules["langchain.docstore.document"] = doc_mod

        if "langchain.docstore" not in sys.modules:
            docstore = types.ModuleType("langchain.docstore")
            docstore.document = sys.modules["langchain.docstore.document"]
            sys.modules["langchain.docstore"] = docstore

        if "langchain.text_splitter" not in sys.modules:
            split_mod = types.ModuleType("langchain.text_splitter")
            split_mod.RecursiveCharacterTextSplitter = RecursiveCharacterTextSplitter
            sys.modules["langchain.text_splitter"] = split_mod

        import os

        os.environ["PADDLE_PDX_DISABLE_MODEL_SOURCE_CHECK"] = "True"
        from paddleocr import PaddleOCR

        # Initialize PaddleOCR
        _paddle_ocr_instance = PaddleOCR(lang="it")
    return _paddle_ocr_instance


def run_paddle_ocr_on_image(contents: bytes, ext: str) -> str:
    """Run PaddleOCR on image bytes by saving it to a temporary file."""
    import os
    import tempfile

    # Write image bytes to a temp file
    suffix = f".{ext.lower()}"
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as f:
        f.write(contents)
        temp_path = f.name

    try:
        ocr = get_paddle_ocr()
        results = list(ocr.predict(temp_path))
        text_lines = []
        for res in results:
            if isinstance(res, dict) and "rec_texts" in res:
                text_lines.extend(res["rec_texts"])
        return "\n".join(text_lines)
    except Exception as e:
        logger.warning(f"PaddleOCR image parsing failed: {e}")
        return ""
    finally:
        try:
            os.remove(temp_path)
        except Exception:
            pass


def run_paddle_ocr_on_pdf(contents: bytes) -> str:
    """Render PDF pages using PyMuPDF (fitz) and run PaddleOCR on them."""
    try:
        import fitz
    except ImportError:
        logger.warning("PyMuPDF (fitz) not installed — skipping OCR PDF rendering")
        return ""

    try:
        doc = fitz.open(stream=contents, filetype="pdf")
    except Exception as e:
        logger.error(f"Failed to open PDF with fitz: {e}")
        return ""

    full_text = []
    for page_num in range(len(doc)):
        try:
            page = doc.load_page(page_num)
            pix = page.get_pixmap()
            png_bytes = pix.tobytes("png")
            page_text = run_paddle_ocr_on_image(png_bytes, "png")
            if page_text:
                full_text.append(page_text)
        except Exception as e:
            logger.warning(f"Failed to run PaddleOCR on PDF page {page_num}: {e}")

    return "\n\n".join(full_text)


async def run_vision_ocr(contents: bytes, ext: str) -> str:
    """Use VisionService to extract text from an image when local OCR is unavailable."""
    try:
        import base64 as _base64

        from core.services.vision.models import (
            ImageContent,
            VisionCapability,
            VisionRequest,
        )
        from core.services.vision.service import VisionService

        mime = (
            "image/jpeg" if ext.lower() in ("jpg", "jpeg") else f"image/{ext.lower()}"
        )
        b64 = _base64.b64encode(contents).decode()
        image = ImageContent.from_base64(b64, media_type=mime)
        request = VisionRequest(
            prompt="Estrai il testo esatto presente in questa immagine, mantenendo la struttura originale. Restituisci solo il testo estratto, senza commenti o prefissi.",
            images=[image],
            capability=VisionCapability.OCR,
            max_tokens=4000,
        )
        service = VisionService()
        response = await service.analyze(request)
        return response.content or ""
    except Exception as e:
        logger.warning(f"Vision OCR fallback failed: {e}")
        return ""


def parse_docx(contents: bytes) -> str:
    """Parse DOCX file contents and return extracted text."""
    import docx

    try:
        doc = docx.Document(io.BytesIO(contents))
        parts = []
        for para in doc.paragraphs:
            if para.text:
                parts.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                row_texts = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_texts.append(cell.text.strip())
                if row_texts:
                    parts.append(" | ".join(row_texts))
        return "\n".join(parts)
    except Exception as e:
        logger.error(f"Failed to parse DOCX: {e}")
        return ""


def parse_pptx(contents: bytes) -> str:
    """Parse PPTX slide deck contents and return extracted text."""
    from pptx import Presentation

    try:
        prs = Presentation(io.BytesIO(contents))
        parts = []
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                parts.append(f"--- Slide {i + 1} ---\n" + "\n".join(slide_texts))
        return "\n\n".join(parts)
    except Exception as e:
        logger.error(f"Failed to parse PPTX: {e}")
        return ""


def create_router(plugin_instance: Any) -> APIRouter:
    """Create and configure the API router for the study assistant plugin."""
    router = APIRouter()
    agent = StudyAgent()

    # Helper: map professor strictness to OpenAI voice
    def get_professor_voice(strictness: str) -> str:
        s = strictness.lower()
        if s == "amichevole":
            return "nova"  # Friendly, upbeat female
        elif s == "scrupoloso":
            return "onyx"  # Deep, authoritative male
        else:
            return "shimmer"  # Balanced professional female

    @router.get("/subjects")
    async def get_subjects():
        """Get all courses."""
        return await StudyDAO.get_subjects()

    @router.post("/subjects")
    async def create_subject(data: SubjectCreate):
        """Create a new course/subject."""
        subject_id = await StudyDAO.create_subject(data.name, data.description)
        if subject_id == -1:
            raise HTTPException(status_code=500, detail="Failed to create subject")
        return {"id": subject_id, "name": data.name, "description": data.description}

    @router.delete("/subjects/{subject_id}")
    async def delete_subject(subject_id: int):
        """Delete a course and all its related materials."""
        success = await StudyDAO.delete_subject(subject_id)
        if not success:
            raise HTTPException(status_code=404, detail="Subject not found")
        return {"success": True}

    @router.post("/subjects/{subject_id}/upload")
    async def upload_document(
        subject_id: int,
        file: UploadFile = File(...),
        relative_path: str | None = Form(None),
        parent_folder_id: str | None = Form(None),
    ):
        """Upload course material (PDF, TXT, MD) and extract text."""
        subjects = await StudyDAO.get_subjects()
        if not any(s["id"] == subject_id for s in subjects):
            raise HTTPException(status_code=404, detail="Subject not found")

        contents = await file.read()
        raw_text = ""
        file_type = (
            file.filename.split(".")[-1].lower() if "." in file.filename else "txt"
        )

        try:
            if file_type == "pdf":
                # Try standard pypdf extraction first
                import pypdf

                try:
                    pdf_reader = pypdf.PdfReader(io.BytesIO(contents))
                    text_list = []
                    for page in pdf_reader.pages:
                        text = page.extract_text()
                        if text:
                            text_list.append(text)
                    raw_text = "\n".join(text_list)
                except Exception as e:
                    logger.warning(
                        f"Standard PDF text extraction failed for {file.filename}, falling back to OCR: {e}"
                    )
                    raw_text = ""

                # Fallback if extracted text is empty or very short (< 50 chars)
                if len(raw_text.strip()) < 50:
                    logger.info(
                        f"PDF {file.filename} extracted text empty or short ({len(raw_text)} chars). Running PaddleOCR..."
                    )
                    raw_text = run_paddle_ocr_on_pdf(contents)

                # If still empty, try extracting embedded images via pypdf + Vision OCR
                if len(raw_text.strip()) < 50:
                    logger.info(
                        f"Trying pypdf image extraction + Vision OCR for {file.filename}..."
                    )
                    try:
                        pdf_reader2 = pypdf.PdfReader(io.BytesIO(contents))
                        page_texts = []
                        for page in pdf_reader2.pages[:10]:
                            for img_obj in page.images:
                                img_text = await run_vision_ocr(img_obj.data, "jpeg")
                                if img_text.strip():
                                    page_texts.append(img_text)
                        if page_texts:
                            raw_text = "\n\n".join(page_texts)
                    except Exception as e:
                        logger.warning(
                            f"pypdf image extraction + Vision OCR failed: {e}"
                        )

            elif file_type == "docx":
                raw_text = parse_docx(contents)
            elif file_type == "pptx":
                raw_text = parse_pptx(contents)
            elif file_type in ("png", "jpg", "jpeg", "webp"):
                raw_text = run_paddle_ocr_on_image(contents, file_type)
                if not raw_text.strip():
                    logger.info(
                        f"Local OCR unavailable for {file.filename}, trying Vision API..."
                    )
                    raw_text = await run_vision_ocr(contents, file_type)
            elif file_type in ("txt", "md"):
                try:
                    raw_text = contents.decode("utf-8")
                except UnicodeDecodeError:
                    raw_text = contents.decode("latin-1")
            else:
                # Fallback decoding for unspecified text files
                try:
                    raw_text = contents.decode("utf-8")
                except UnicodeDecodeError:
                    raw_text = contents.decode("latin-1")
        except Exception as e:
            logger.error(f"Error parsing file {file.filename}: {e}")
            raise HTTPException(
                status_code=400, detail=f"Failed to parse file: {e}"
            ) from e

        if not raw_text.strip():
            raw_text = f"[{file.filename}]\n[Contenuto non estraibile automaticamente — documento salvato senza testo]"

        # Resolve parent_folder_id from form payload
        pf_id = None
        if parent_folder_id and parent_folder_id not in (
            "null",
            "undefined",
            "None",
            "",
        ):
            try:
                pf_id = int(parent_folder_id)
            except ValueError:
                pass

        # Resolve folder_id from relative_path and pf_id
        folder_id = pf_id
        if relative_path:
            parts = [p.strip() for p in relative_path.split("/") if p.strip()]
            # If relative path ends with file.filename, remove it
            if parts and parts[-1] == file.filename:
                parts.pop()

            curr_parent_id = pf_id
            for part in parts:
                all_folders = await StudyDAO.get_folders(subject_id)
                matched = None
                for f in all_folders:
                    p_id = f.get("parent_id")
                    is_root_match = (curr_parent_id is None) and p_id is None
                    is_sub_match = curr_parent_id is not None and p_id == curr_parent_id
                    if f["name"].lower() == part.lower() and (
                        is_root_match or is_sub_match
                    ):
                        matched = f
                        break

                if matched:
                    curr_parent_id = matched["id"]
                else:
                    new_fid = await StudyDAO.create_folder(
                        subject_id, part, curr_parent_id
                    )
                    if new_fid != -1:
                        curr_parent_id = new_fid
                    else:
                        # Race condition: another parallel upload already created this folder — re-look it up
                        retry_folders = await StudyDAO.get_folders(subject_id)
                        for f in retry_folders:
                            p_id = f.get("parent_id")
                            is_root_match = (curr_parent_id is None) and p_id is None
                            is_sub_match = (
                                curr_parent_id is not None and p_id == curr_parent_id
                            )
                            if f["name"].lower() == part.lower() and (
                                is_root_match or is_sub_match
                            ):
                                curr_parent_id = f["id"]
                                break

            folder_id = curr_parent_id

        doc_id = await StudyDAO.add_document(
            subject_id=subject_id,
            name=file.filename,
            file_path="",  # Local path is empty since we store raw text
            file_type=file_type,
            raw_text=raw_text,
            folder_id=folder_id,
        )

        if doc_id == -1:
            raise HTTPException(status_code=500, detail="Failed to save document")

        return {
            "id": doc_id,
            "name": file.filename,
            "character_count": len(raw_text),
            "folder_id": folder_id,
        }

    @router.get("/subjects/{subject_id}/documents")
    async def get_documents(subject_id: int):
        """Get list of uploaded materials."""
        return await StudyDAO.get_documents(subject_id)

    @router.delete("/documents/{doc_id}")
    async def delete_document(doc_id: int):
        """Delete uploaded document."""
        success = await StudyDAO.delete_document(doc_id)
        if not success:
            raise HTTPException(status_code=404, detail="Document not found")
        return {"success": True}

    @router.get("/subjects/{subject_id}/folders")
    async def get_folders(subject_id: int):
        """Get all folders in the subject."""
        return await StudyDAO.get_folders(subject_id)

    @router.post("/subjects/{subject_id}/folders")
    async def create_folder(
        subject_id: int, name: str = Form(...), parent_id: str | None = Form(None)
    ):
        """Create a new folder."""
        p_id = None
        if parent_id and parent_id not in ("null", "undefined", "None", ""):
            try:
                p_id = int(parent_id)
            except ValueError:
                pass

        folder_id = await StudyDAO.create_folder(subject_id, name, p_id)
        if folder_id == -1:
            raise HTTPException(status_code=500, detail="Failed to create folder")
        return {"id": folder_id, "name": name, "parent_id": p_id}

    @router.delete("/folders/{folder_id}")
    async def delete_folder(folder_id: int):
        """Delete a folder and all its contents (cascading)."""
        success = await StudyDAO.delete_folder(folder_id)
        if not success:
            raise HTTPException(status_code=404, detail="Folder not found")
        return {"success": True}

    @router.put("/folders/{folder_id}")
    async def rename_folder(folder_id: int, data: dict[str, str]):
        """Rename a folder."""
        name = data.get("name")
        if not name:
            raise HTTPException(status_code=400, detail="Name is required")
        success = await StudyDAO.rename_folder(folder_id, name)
        if not success:
            raise HTTPException(status_code=404, detail="Folder not found")
        return {"success": True}

    @router.put("/documents/{doc_id}")
    async def update_document(doc_id: int, data: dict[str, Any]):
        """Rename or move a document."""
        # Check if we are renaming
        if "name" in data:
            name = data["name"]
            success = await StudyDAO.rename_document(doc_id, name)
            if not success:
                raise HTTPException(status_code=404, detail="Document not found")

        # Check if we are moving
        if "folder_id" in data:
            folder_id_val = data["folder_id"]
            f_id = None
            if folder_id_val and str(folder_id_val) not in (
                "null",
                "undefined",
                "None",
                "",
            ):
                try:
                    f_id = int(folder_id_val)
                except ValueError:
                    pass
            success = await StudyDAO.move_document(doc_id, f_id)
            if not success:
                raise HTTPException(status_code=404, detail="Document not found")

        return {"success": True}

    @router.put("/folders/{folder_id}/move")
    async def move_folder(folder_id: int, data: dict[str, Any]):
        """Move a folder under a new parent."""
        parent_id_val = data.get("parent_id")
        p_id = None
        if parent_id_val and str(parent_id_val) not in (
            "null",
            "undefined",
            "None",
            "",
        ):
            try:
                p_id = int(parent_id_val)
            except ValueError:
                pass
        success = await StudyDAO.move_folder(folder_id, p_id)
        if not success:
            raise HTTPException(
                status_code=400,
                detail="Failed to move folder (check for cycles or same parent)",
            )
        return {"success": True}

    @router.post("/subjects/{subject_id}/bulk-delete")
    async def bulk_delete(subject_id: int, req: BulkDeleteRequest):
        """Delete multiple documents and folders in bulk, scoped to this subject."""
        valid_doc_ids = {d["id"] for d in await StudyDAO.get_documents(subject_id)}
        valid_folder_ids = {f["id"] for f in await StudyDAO.get_folders(subject_id)}

        deleted_docs = 0
        deleted_folders = 0
        for doc_id in req.document_ids:
            if doc_id in valid_doc_ids and await StudyDAO.delete_document(doc_id):
                deleted_docs += 1
        for folder_id in req.folder_ids:
            if folder_id in valid_folder_ids and await StudyDAO.delete_folder(
                folder_id
            ):
                deleted_folders += 1
        return {
            "success": True,
            "deleted_documents": deleted_docs,
            "deleted_folders": deleted_folders,
        }

    @router.post("/subjects/{subject_id}/bulk-move")
    async def bulk_move(subject_id: int, req: BulkMoveRequest):
        """Move multiple documents and folders in bulk to target_folder_id."""
        target_id = req.target_folder_id
        if target_id and str(target_id) in ("null", "undefined", "None", ""):
            target_id = None

        if target_id is not None:
            folders = await StudyDAO.get_folders(subject_id)
            if not any(f["id"] == target_id for f in folders):
                raise HTTPException(
                    status_code=400, detail="Target folder not found in this subject"
                )

        moved_docs = 0
        moved_folders = 0
        failed_folders = []

        for folder_id in req.folder_ids:
            if folder_id == target_id:
                failed_folders.append(folder_id)
                continue
            success = await StudyDAO.move_folder(folder_id, target_id)
            if success:
                moved_folders += 1
            else:
                failed_folders.append(folder_id)

        for doc_id in req.document_ids:
            success = await StudyDAO.move_document(doc_id, target_id)
            if success:
                moved_docs += 1

        if failed_folders:
            return {
                "success": True,
                "moved_documents": moved_docs,
                "moved_folders": moved_folders,
                "failed_folders": failed_folders,
                "message": f"Moved {moved_docs} files and {moved_folders} folders. Failed for folders: {failed_folders} (cycles).",
            }

        return {
            "success": True,
            "moved_documents": moved_docs,
            "moved_folders": moved_folders,
        }

    @router.get("/subjects/{subject_id}/flashcards")
    async def get_flashcards(subject_id: int, due_only: bool = False):
        """Get flashcard deck, optionally filtered for cards due to review."""
        cards = await StudyDAO.get_flashcards(subject_id, due_only=due_only)
        stats = await StudyDAO.get_flashcard_count_stats(subject_id)
        return {"flashcards": cards, "stats": stats}

    @router.post("/subjects/{subject_id}/flashcards/generate")
    async def generate_flashcards(subject_id: int, req: FlashcardGenerateRequest):
        """Extract flashcards automatically from subject notes."""
        cards = await agent.generate_flashcards(subject_id, req.num_flashcards)
        if not cards:
            raise HTTPException(
                status_code=400,
                detail="No documents found or failed to parse content for flashcard generation",
            )

        # Save to DB
        db_cards = [(c["question"], c["answer"]) for c in cards]
        await StudyDAO.create_flashcards(subject_id, db_cards)

        return {"success": True, "count": len(cards)}

    @router.post("/flashcards/review")
    async def review_flashcard(req: FlashcardReview):
        """Submit flashcard review score (0-5) and reschedule it using SM-2."""
        try:
            res = await StudyDAO.update_flashcard_sm2(req.flashcard_id, req.rating)
            return res
        except ValueError as e:
            raise HTTPException(status_code=404, detail=str(e)) from e

    @router.post("/sessions/oral/start")
    async def start_oral_session(req: OralSessionStart):
        """Initiate a simulated oral interrogation."""
        # 1. Fetch documents to extract syllabus topics
        docs = await StudyDAO.get_documents(req.subject_id)

        # Determine topics
        topics = []
        if docs:
            # Construct a prompt to extract 5 topics
            combined_titles = ", ".join([d["name"] for d in docs])
            prompt = f"Basandoti sui seguenti file caricati per un corso universitario: {combined_titles}. Identifica 5 argomenti o capitoli principali su cui fare domande all'esame. Restituisci gli argomenti come una lista separata da virgole, senza numeri o introduzione."
            try:
                llm = get_llm_service()
                res = await llm.generate_response(prompt)
                topics = [t.strip() for t in res.split(",") if t.strip()]
            except Exception as e:
                logger.error(f"Failed to generate topics from files: {e}")

        # Fallback topics if none extracted
        if not topics:
            subjects = await StudyDAO.get_subjects()
            subject = next((s for s in subjects if s["id"] == req.subject_id), None)
            subject_name = subject["name"] if subject else "Generale"
            prompt = f"Genera 5 argomenti accademici principali per la materia universitaria '{subject_name}'. Restituisci i capitoli separati solo da virgole."
            try:
                llm = get_llm_service()
                res = await llm.generate_response(prompt)
                topics = [t.strip() for t in res.split(",") if t.strip()]
            except Exception:
                topics = [
                    "Capitolo 1",
                    "Capitolo 2",
                    "Capitolo 3",
                    "Capitolo 4",
                    "Capitolo 5",
                ]

        # Limit to 5 topics
        topics = topics[:5]

        session_id = await StudyDAO.create_session(
            subject_id=req.subject_id,
            professor_name=req.professor_name,
            strictness=req.strictness,
            difficulty_level=req.difficulty_level,
        )

        if session_id == -1:
            raise HTTPException(status_code=500, detail="Failed to create exam session")

        # Save topics list inside session metadata / first transcript block
        session_info = await StudyDAO.get_session(session_id)
        transcript = session_info["transcript"] if session_info else []
        transcript.append(
            {
                "type": "system",
                "topics": topics,
                "message": f"Esame orale avviato con {req.professor_name}.",
            }
        )
        await StudyDAO.update_session(session_id, None, "active", transcript, None)

        return {
            "session_id": session_id,
            "professor_name": req.professor_name,
            "strictness": req.strictness,
            "difficulty_level": req.difficulty_level,
            "topics": topics,
        }

    @router.post("/sessions/oral/next")
    async def next_oral_question(session_id: int = Form(...)):
        """Choose and formulate the next question using MCTS and speak it via OpenAI TTS."""
        session = await StudyDAO.get_session(session_id)
        if not session:
            raise HTTPException(status_code=404, detail="Session not found")
        if session["status"] == "completed":
            return {"status": "completed"}

        # Extract topics from first transcript block
        transcript = session["transcript"]
        topics = []
        for block in transcript:
            if block.get("type") == "system" and "topics" in block:
                topics = block["topics"]
                break

        if not topics:
            topics = ["Argomenti Generali"]

        # Build scores and past questions list
        current_scores = {}
        question_history = []
        for block in transcript:
            if block.get("type") == "evaluation":
                topic = block.get("topic")
                score = block.get("score", 6.0)
                current_scores[topic] = score
            if block.get("type") == "question":
                question_history.append(
                    {"topic": block.get("topic"), "style": block.get("style")}
                )

        # Check if we should end the session
        # If we asked about 4 topics, or 5 questions in total, finalize
        if len(current_scores) >= 4 or len(question_history) >= 5:
            # End of interrogation! Calculate final grade
            avg_score = (
                sum(current_scores.values()) / len(current_scores)
                if current_scores
                else 6.0
            )

            # Map 0-10 score to Italian 18-30 scale
            if avg_score < 6.0:
                final_grade = "Bocciato (Insufficiente)"
            else:
                grade_num = round(18 + (avg_score - 6.0) * 3)
                final_grade = f"{grade_num}/30"
                if grade_num > 30:
                    final_grade = "30 e Lode"

            transcript.append(
                {
                    "type": "system_finish",
                    "avg_score": avg_score,
                    "final_grade": final_grade,
                    "message": f"Esame terminato. Voto finale: {final_grade}.",
                }
            )
            await StudyDAO.update_session(
                session_id, avg_score, "completed", transcript, None
            )
            return {
                "status": "completed",
                "avg_score": avg_score,
                "final_grade": final_grade,
                "transcript": transcript,
            }

        # Select next question parameters using MCTS
        planner = OralExamPlanner(
            topics=topics,
            strictness=session["strictness"],
            difficulty_level=session["difficulty_level"],
        )

        next_params = await planner.select_next_question(
            current_scores=current_scores,
            question_history=question_history,
            current_topic=session["current_topic"],
        )

        topic = next_params.get("topic")
        style = next_params.get("style")

        # Get documents contents for RAG context
        docs = await StudyDAO.get_documents(session["subject_id"])
        full_docs = []
        for d in docs:
            fd = await StudyDAO.get_document(d["id"])
            if fd:
                full_docs.append(fd)

        # Formulate question text via LLM
        question_text = await agent.formulate_oral_question(
            subject_name=session["subject_name"],
            topic=topic,
            style=style,
            strictness=session["strictness"],
            difficulty=session["difficulty_level"],
            context_docs=full_docs,
        )

        # Text-To-Speech generation
        voice = get_professor_voice(session["strictness"])
        audio_base64 = None
        try:
            voice_service = VoiceService()
            voice_res = await voice_service.text_to_speech(
                text=question_text,
                voice=voice,
                format=AudioFormat.MP3,
                provider=VoiceProvider.OPENAI,
            )
            audio_base64 = voice_res.audio_base64
        except Exception as e:
            logger.error(f"TTS generation failed: {e}")

        # Update session with new question
        transcript.append(
            {
                "type": "question",
                "topic": topic,
                "style": style,
                "text": question_text,
            }
        )
        await StudyDAO.update_session(session_id, None, "active", transcript, topic)

        return {
            "status": "active",
            "topic": topic,
            "style": style,
            "text": question_text,
            "audio": audio_base64,
        }

    @router.post("/sessions/oral/answer")
    async def answer_oral_question(data: OralSessionAnswer):
        """Submit answer to the current question, evaluate, and return feedback."""
        session = await StudyDAO.get_session(data.session_id)
        if not session:
            raise HTTPException(status_code=404, detail="Session not found")
        if session["status"] == "completed":
            return {"status": "completed"}

        transcript = session["transcript"]
        current_topic = session["current_topic"]

        # Find the last question asked
        last_question = None
        for block in reversed(transcript):
            if block.get("type") == "question":
                last_question = block
                break

        if not last_question:
            raise HTTPException(
                status_code=400, detail="No active question found to answer"
            )

        # Evaluate answer
        eval_result = await agent.evaluate_oral_answer(
            subject_name=session["subject_name"],
            topic=current_topic,
            question=last_question["text"],
            answer=data.answer_text,
            strictness=session["strictness"],
            difficulty=session["difficulty_level"],
        )

        # Update transcript
        transcript.append(
            {
                "type": "answer",
                "text": data.answer_text,
            }
        )
        transcript.append(
            {
                "type": "evaluation",
                "topic": current_topic,
                "score": eval_result["score"],
                "feedback": eval_result["feedback"],
                "is_correct": eval_result["is_correct"],
            }
        )

        await StudyDAO.update_session(
            session_id=data.session_id,
            score=None,
            status="active",
            transcript=transcript,
            current_topic=current_topic,
        )

        return {
            "score": eval_result["score"],
            "feedback": eval_result["feedback"],
            "is_correct": eval_result["is_correct"],
        }

    @router.post("/sessions/oral/transcribe")
    async def transcribe_audio(file: UploadFile = File(...)):
        """Transcribe oral student response audio using OpenAI Whisper."""
        try:
            audio_bytes = await file.read()
            voice_service = VoiceService()
            stt_res = await voice_service.speech_to_text(
                audio_data=audio_bytes, language="it", provider=VoiceProvider.OPENAI
            )
            return {"text": stt_res.text}
        except Exception as e:
            logger.error(f"STT transcription failed: {e}")
            raise HTTPException(
                status_code=500, detail=f"Transcription failed: {e}"
            ) from e

    @router.post("/sessions/oral/speak")
    async def speak_text(text: str = Form(...), strictness: str = Form("equo")):
        """Synthesize arbitrary text (e.g. general chat) using OpenAI TTS."""
        voice = get_professor_voice(strictness)
        try:
            voice_service = VoiceService()
            voice_res = await voice_service.text_to_speech(
                text=text,
                voice=voice,
                format=AudioFormat.MP3,
                provider=VoiceProvider.OPENAI,
            )
            return {"audio": voice_res.audio_base64}
        except Exception as e:
            logger.error(f"TTS speak failed: {e}")
            raise HTTPException(status_code=500, detail=f"TTS failed: {e}") from e

    @router.post("/subjects/{subject_id}/chat")
    async def chat_with_assistant(subject_id: int, query: str = Form(...)):
        """General tutoring chat with RAG documents."""
        res = await agent.answer_study_query(subject_id, query)
        return res

    @router.get("/active/feynman/concept")
    async def get_feynman_concept(subject_id: int):
        """Extract or suggest a concept to study using Feynman technique."""
        concept = await agent.select_feynman_concept(subject_id)
        return {"concept": concept}

    @router.post("/active/feynman/respond")
    async def respond_feynman_step(
        subject_id: int = Form(...),
        concept_name: str = Form(...),
        explanation: str = Form(...),
        history_json: str = Form("[]"),
    ):
        """Evaluate a Feynman explanation step and return feedback/next question."""
        try:
            history = json.loads(history_json)
        except Exception:
            history = []
        res = await agent.evaluate_feynman_step(
            subject_id=subject_id,
            concept_name=concept_name,
            explanation=explanation,
            history=history,
        )
        return res

    @router.post("/subjects/{subject_id}/deconstruct")
    async def deconstruct_subject_concepts(
        subject_id: int, force_refresh: bool = False
    ):
        """Deconstruct course materials into focal cheat-sheets, probable questions, hooks, and conceptual maps."""
        res = await agent.deconstruct_subject_concepts(
            subject_id, force_refresh=force_refresh
        )
        return res

    @router.get("/subjects/{subject_id}/topics")
    async def get_subject_topics(subject_id: int):
        """Extract study topics directly from subject documents using LLM."""
        topics = await agent.extract_document_topics(subject_id)
        return {"topics": topics}

    @router.post("/subjects/{subject_id}/podcasts/generate-stream")
    async def generate_podcast_stream(subject_id: int, req: PodcastGenerateRequest):
        """Generate a podcast and stream progress events."""
        import asyncio
        import re

        from fastapi.responses import StreamingResponse

        async def event_generator():
            import os

            # 1. Fetch subject details
            subjects = await StudyDAO.get_subjects()
            subject = next((s for s in subjects if s["id"] == subject_id), None)
            subject_name = subject["name"] if subject else "Generale"

            yield (
                json.dumps(
                    {
                        "status": "analyzing",
                        "message": "Analisi dei materiali del corso...",
                    }
                )
                + "\n"
            )
            await asyncio.sleep(0.01)

            # 2. Choose professor profile
            try:
                prof_profile = await agent.choose_podcast_professor(
                    subject_name, req.topic
                )
                professor_name = prof_profile.get(
                    "professor_name", "Prof. Alessandro Neri"
                )
                voice = prof_profile.get("voice", "echo")
            except Exception:
                professor_name = "Prof. Alessandro Neri"
                voice = "echo"

            # 3. Build RAG context
            docs = await StudyDAO.get_documents(subject_id)
            context_parts = []
            for doc in docs:
                full_doc = await StudyDAO.get_document(doc["id"])
                if full_doc and full_doc.get("raw_text"):
                    context_parts.append(f"[{doc['name']}]\n{full_doc['raw_text']}")
                    if sum(len(p) for p in context_parts) >= 14000:
                        break
            context = (
                "\n\n---\n\n".join(context_parts)[:14000]
                or "Usa le tue conoscenze generali."
            )

            # 4. Determine number of episodes from depth level
            if req.depth_level == "breve":
                num_episodes = 1
            elif req.depth_level == "approfondito":
                num_episodes = 3
            else:
                num_episodes = 2

            # 5. STEP 1 — generate podcast outline (title + episode titles only, short response)
            yield (
                json.dumps(
                    {
                        "status": "scripting",
                        "message": "Generazione struttura del podcast...",
                    }
                )
                + "\n"
            )
            await asyncio.sleep(0.01)

            outline_prompt = f"""Sei {professor_name}, docente universitario. Stai pianificando un podcast didattico di ESATTAMENTE {num_episodes} puntate sulla materia '{subject_name}', argomento: '{req.topic}'.

Restituisci un oggetto JSON con questa struttura ESATTA (NIENTE altro testo, niente markdown):
{{
  "title": "Titolo generale coinvolgente del podcast",
  "episodes": [
    {{"episode_number": 1, "title": "Titolo puntata 1 — aspetto specifico dell'argomento"}},
    {{"episode_number": 2, "title": "Titolo puntata 2 — aspetto specifico dell'argomento"}}
  ]
}}

L'array "episodes" deve contenere ESATTAMENTE {num_episodes} elementi. Ogni titolo deve coprire un aspetto diverso e complementare dell'argomento.
Rispondi SOLO con il JSON."""

            outline_data = {}
            try:
                outline_resp = await agent.llm_service.generate_response(outline_prompt)
                cleaned = outline_resp.strip()
                if cleaned.startswith("```"):
                    cleaned = re.sub(r"^```(json)?\n", "", cleaned)
                    cleaned = re.sub(r"\n```$", "", cleaned)
                outline_data = json.loads(cleaned.strip())
            except Exception as e:
                logger.error(f"Failed to generate podcast outline: {e}")
                yield (
                    json.dumps(
                        {
                            "status": "error",
                            "message": f"Generazione struttura fallita: {e}",
                        }
                    )
                    + "\n"
                )
                return

            if not outline_data.get("title") or not outline_data.get("episodes"):
                yield (
                    json.dumps(
                        {
                            "status": "error",
                            "message": "L'LLM non ha generato una struttura valida.",
                        }
                    )
                    + "\n"
                )
                return

            podcast_title = outline_data["title"]
            raw_outlines = outline_data["episodes"]
            # Normalize: ensure sequential episode_number starting at 1, cap to num_episodes
            episode_outlines = []
            for i, ep in enumerate(raw_outlines[:num_episodes]):
                episode_outlines.append(
                    {
                        "episode_number": i + 1,
                        "title": ep.get("title", f"Puntata {i + 1}"),
                    }
                )
            # If LLM returned fewer episodes than requested, fill remaining
            while len(episode_outlines) < num_episodes:
                n = len(episode_outlines) + 1
                episode_outlines.append({"episode_number": n, "title": f"Puntata {n}"})

            # 6. STEP 2 — generate each episode script in parallel (one LLM call per episode)
            yield (
                json.dumps(
                    {
                        "status": "scripting",
                        "message": f"Generazione di {len(episode_outlines)} script in parallelo...",
                        "total_episodes": len(episode_outlines),
                        "episode_outlines": [
                            {
                                "episode_number": ep["episode_number"],
                                "title": ep["title"],
                            }
                            for ep in episode_outlines
                        ],
                    }
                )
                + "\n"
            )
            await asyncio.sleep(0.01)

            async def gen_script(ep_info: dict) -> dict:
                try:
                    script = await agent.generate_episode_script(
                        professor_name=professor_name,
                        subject_name=subject_name,
                        topic=req.topic,
                        episode_number=ep_info["episode_number"],
                        episode_title=ep_info["title"],
                        num_episodes=len(episode_outlines),
                        context=context,
                    )
                    return {
                        "episode_number": ep_info["episode_number"],
                        "title": ep_info["title"],
                        "script_text": script.strip(),
                    }
                except Exception as e:
                    logger.error(
                        f"Script gen failed for ep {ep_info['episode_number']}: {e}"
                    )
                    # Return a minimal placeholder so the episode slot is preserved
                    return {
                        "episode_number": ep_info["episode_number"],
                        "title": ep_info["title"],
                        "script_text": f"Script per la puntata {ep_info['episode_number']} non disponibile.",
                    }

            episodes_with_scripts = await asyncio.gather(
                *[gen_script(ep) for ep in episode_outlines]
            )

            # 7. Create podcast DB record
            podcast_id = await StudyDAO.create_podcast(
                subject_id=subject_id,
                title=podcast_title,
                topic=req.topic,
                professor_voice=voice,
                professor_name=professor_name,
                depth_level=req.depth_level,
            )
            if podcast_id == -1:
                yield (
                    json.dumps(
                        {
                            "status": "error",
                            "message": "Salvataggio metadati podcast fallito.",
                        }
                    )
                    + "\n"
                )
                return

            plugin_dir = os.path.dirname(os.path.abspath(__file__))
            podcasts_dir = os.path.join(plugin_dir, "static", "podcasts")
            os.makedirs(podcasts_dir, exist_ok=True)

            voice_service = VoiceService()

            yield (
                json.dumps(
                    {
                        "status": "synthesizing",
                        "message": f"Sintesi vocale di {len(episodes_with_scripts)} puntate...",
                        "total_episodes": len(episodes_with_scripts),
                        "completed_episodes": 0,
                    }
                )
                + "\n"
            )
            await asyncio.sleep(0.01)

            # 8. Synthesize in parallel, yield progress per episode
            # asyncio.wait returns the original futures (unlike as_completed which wraps them)
            completed_count = 0
            generated_episodes = [None] * len(episodes_with_scripts)
            future_to_index = {}
            for idx, ep in enumerate(episodes_with_scripts):
                fut = asyncio.ensure_future(
                    agent._synthesize_episode(
                        ep=ep,
                        podcast_id=podcast_id,
                        voice=voice,
                        voice_service=voice_service,
                        podcasts_dir=podcasts_dir,
                    )
                )
                future_to_index[fut] = idx

            pending = set(future_to_index.keys())
            while pending:
                done, pending = await asyncio.wait(
                    pending, return_when=asyncio.FIRST_COMPLETED
                )
                for fut in done:
                    idx = future_to_index[fut]
                    try:
                        ep_res = fut.result()
                        generated_episodes[idx] = ep_res
                        completed_count += 1
                        yield (
                            json.dumps(
                                {
                                    "status": "episode_ready",
                                    "message": f"Puntata {ep_res['episode_number']}/{len(episodes_with_scripts)} pronta: '{ep_res['title']}'",
                                    "episode_number": ep_res["episode_number"],
                                    "total_episodes": len(episodes_with_scripts),
                                    "completed_episodes": completed_count,
                                }
                            )
                            + "\n"
                        )
                    except Exception as e:
                        logger.error(
                            f"TTS synthesis failed for episode at index {idx}: {e}"
                        )
                        completed_count += 1
                        yield (
                            json.dumps(
                                {
                                    "status": "episode_error",
                                    "message": f"Sintesi puntata {idx + 1} fallita: {e}",
                                    "episode_number": idx + 1,
                                    "total_episodes": len(episodes_with_scripts),
                                    "completed_episodes": completed_count,
                                }
                            )
                            + "\n"
                        )
                    await asyncio.sleep(0.01)

            successful_episodes = [e for e in generated_episodes if e is not None]

            if not successful_episodes:
                # Every episode failed synthesis — don't leave an empty podcast behind.
                try:
                    await StudyDAO.delete_podcast(podcast_id)
                except Exception as cleanup_err:
                    logger.error(
                        f"Failed to clean up orphaned podcast {podcast_id}: {cleanup_err}"
                    )
                yield (
                    json.dumps(
                        {
                            "status": "error",
                            "message": "Sintesi vocale fallita per tutte le puntate. Podcast non salvato.",
                        }
                    )
                    + "\n"
                )
                return

            partial_failure = len(successful_episodes) < len(episodes_with_scripts)
            yield (
                json.dumps(
                    {
                        "status": "completed_with_errors"
                        if partial_failure
                        else "completed",
                        "message": (
                            f"Podcast generato con {len(successful_episodes)}/{len(episodes_with_scripts)} puntate "
                            "(alcune puntate non sono state sintetizzate)."
                            if partial_failure
                            else "Podcast generato con successo!"
                        ),
                        "podcast": {
                            "id": podcast_id,
                            "title": podcast_title,
                            "topic": req.topic,
                            "professor_name": professor_name,
                            "professor_voice": voice,
                            "depth_level": req.depth_level,
                            "episodes": successful_episodes,
                        },
                    }
                )
                + "\n"
            )

        return StreamingResponse(event_generator(), media_type="application/x-ndjson")

    @router.post("/subjects/{subject_id}/podcasts/generate")
    async def generate_podcast(subject_id: int, req: PodcastGenerateRequest):
        """Generate a multi-episode podcast for the subject based on topic and depth."""
        try:
            res = await agent.generate_podcast(
                subject_id=subject_id, topic=req.topic, depth_level=req.depth_level
            )
            return res
        except ValueError as e:
            raise HTTPException(status_code=400, detail=str(e)) from e
        except Exception as e:
            logger.error(f"Failed to generate podcast: {e}")
            raise HTTPException(
                status_code=500,
                detail=f"Errore interno nella generazione del podcast: {e}",
            ) from e

    @router.get("/subjects/{subject_id}/sessions")
    async def get_sessions(subject_id: int):
        """Get all oral exam sessions for a subject."""
        return await StudyDAO.get_sessions(subject_id)

    @router.get("/sessions/oral/{session_id}")
    async def get_oral_session(session_id: int):
        """Get a specific oral exam session by ID with its transcript."""
        session = await StudyDAO.get_session(session_id)
        if not session:
            raise HTTPException(status_code=404, detail="Session non trovata")
        session["session_id"] = session["id"]
        return session

    @router.get("/subjects/{subject_id}/podcasts")
    async def get_podcasts(subject_id: int):
        """Get all podcasts generated for a subject."""
        podcasts = await StudyDAO.get_podcasts(subject_id)
        res = []
        for p in podcasts:
            episodes = await StudyDAO.get_podcast_episodes(p["id"])
            res.append({**p, "episodes": episodes})
        return res

    @router.delete("/podcasts/{podcast_id}")
    async def delete_podcast(podcast_id: int):
        """Delete a podcast and remove all its audio files from disk."""
        import os

        episodes = await StudyDAO.get_podcast_episodes(podcast_id)

        # Get path to static folder
        plugin_dir = os.path.dirname(os.path.abspath(__file__))
        podcasts_dir = os.path.join(plugin_dir, "static", "podcasts")

        # Delete audio files
        for ep in episodes:
            audio_filename = ep.get("audio_filename")
            if audio_filename:
                file_path = os.path.join(podcasts_dir, audio_filename)
                if os.path.exists(file_path):
                    try:
                        os.remove(file_path)
                    except Exception as e:
                        logger.error(f"Failed to delete audio file {file_path}: {e}")

        success = await StudyDAO.delete_podcast(podcast_id)
        if not success:
            raise HTTPException(
                status_code=404, detail="Podcast non trovato o già eliminato"
            )
        return {"success": True}

    @router.get("/debug/events")
    async def get_debug_events():
        """Get all tracked MCTS and Graph RAG debug events."""
        from .debug_tracker import DebugTracker

        return DebugTracker.get_events()

    @router.post("/debug/clear")
    async def clear_debug_events():
        """Clear all debug events."""
        from .debug_tracker import DebugTracker

        DebugTracker.clear()
        return {"success": True}

    return router
